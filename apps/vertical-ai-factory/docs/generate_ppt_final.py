"""
PPT 최종본 - 도형, 다이어그램, 시각적 요소 추가
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 파일 경로
INPUT_FILE = '/home/roarm_m3/ai-factory-lab/apps/vertical-ai-factory/docs/생산계획시스템_설명자료_260108.pptx'
OUTPUT_FILE = '/home/roarm_m3/ai-factory-lab/apps/vertical-ai-factory/docs/생산계획시스템_최종본.pptx'

# 색상 팔레트
COLORS = {
    "navy": RGBColor(0x1E, 0x3A, 0x5F),
    "blue": RGBColor(0x3B, 0x82, 0xF6),
    "light_blue": RGBColor(0x93, 0xC5, 0xFD),
    "orange": RGBColor(0xF5, 0x9E, 0x0B),
    "green": RGBColor(0x10, 0xB9, 0x81),
    "gray": RGBColor(0x6B, 0x72, 0x80),
    "light_gray": RGBColor(0xE5, 0xE7, 0xEB),
    "dark": RGBColor(0x1F, 0x29, 0x37),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}


def add_rounded_rect(slide, left, top, width, height, fill_color, text="", font_size=Pt(12), font_color=None):
    """둥근 사각형 추가"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # 테두리 없음
    
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = font_color or COLORS["white"]
        p.font.bold = True
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].space_before = Pt(0)
        shape.text_frame.paragraphs[0].space_after = Pt(0)
    
    return shape


def add_arrow(slide, left, top, width, height):
    """화살표 추가"""
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["gray"]
    shape.line.fill.background()
    return shape


def create_slide_1_title(slide):
    """슬라이드 1: 표지"""
    # 상단 장식 바
    add_rounded_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.3), COLORS["navy"])
    
    # 하단 장식 바
    add_rounded_rect(slide, Inches(0), Inches(7.2), Inches(10), Inches(0.3), COLORS["blue"])
    
    # 메인 제목 영역
    for shape in slide.shapes:
        if hasattr(shape, 'text') and '차세대' in shape.text:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = "차세대 생산관리 시스템"
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = COLORS["navy"]
            p.alignment = PP_ALIGN.CENTER
            
            p = tf.add_paragraph()
            p.text = "디지털 전환을 통한 제조 혁신"
            p.font.size = Pt(24)
            p.font.color.rgb = COLORS["gray"]
            p.alignment = PP_ALIGN.CENTER


def create_slide_2_toc(slide):
    """슬라이드 2: 목차"""
    # 기존 텍스트 업데이트
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if 'Contents' in shape.text or '목차' in shape.text:
                shape.text_frame.paragraphs[0].text = "목차"
                shape.text_frame.paragraphs[0].font.size = Pt(36)
                shape.text_frame.paragraphs[0].font.color.rgb = COLORS["navy"]
                shape.text_frame.paragraphs[0].font.bold = True


def create_mes_diagram(slide):
    """MES 구성도 다이어그램"""
    start_x = Inches(1.5)
    start_y = Inches(2)
    box_width = Inches(7)
    box_height = Inches(1.2)
    gap = Inches(0.3)
    
    layers = [
        ("ERP 계층", "전사 계획 / 영업 / 구매 / 재무", COLORS["navy"]),
        ("MES 계층", "생산계획 / 실적관리 / 품질 / 설비", COLORS["blue"]),
        ("설비 계층", "PLC / 센서 / SCADA / IoT", COLORS["gray"]),
    ]
    
    for i, (title, desc, color) in enumerate(layers):
        y = start_y + i * (box_height + gap + Inches(0.3))
        
        # 박스
        box = add_rounded_rect(slide, start_x, y, box_width, box_height, color)
        
        # 텍스트 추가
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["light_gray"]
        p.alignment = PP_ALIGN.CENTER
        
        # 화살표 (마지막 제외)
        if i < len(layers) - 1:
            arrow_x = start_x + box_width / 2 - Inches(0.2)
            arrow_y = y + box_height
            add_arrow(slide, arrow_x, arrow_y, Inches(0.4), Inches(0.3))


def create_roadmap_diagram(slide):
    """로드맵 다이어그램"""
    phases = [
        ("1단계", "생산계획", "주간/일간 계획\nOR-Tools 최적화", COLORS["blue"], "현재"),
        ("2단계", "실적관리", "바코드 수집\n실시간 현황", COLORS["light_blue"], ""),
        ("3단계", "원가관리", "원가 집계\n시뮬레이션", COLORS["light_gray"], ""),
    ]
    
    start_x = Inches(0.8)
    y = Inches(2.5)
    box_width = Inches(2.8)
    box_height = Inches(2.5)
    gap = Inches(0.3)
    
    for i, (phase, name, desc, color, badge) in enumerate(phases):
        x = start_x + i * (box_width + gap)
        
        # 메인 박스
        box = add_rounded_rect(slide, x, y, box_width, box_height, color)
        
        tf = box.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["white"] if color != COLORS["light_gray"] else COLORS["dark"]
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = name
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"] if color != COLORS["light_gray"] else COLORS["dark"]
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = ""
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS["white"] if color != COLORS["light_gray"] else COLORS["gray"]
        p.alignment = PP_ALIGN.CENTER
        
        # 현재 뱃지
        if badge:
            badge_shape = add_rounded_rect(slide, x + box_width - Inches(0.8), y - Inches(0.15), 
                                          Inches(0.8), Inches(0.3), COLORS["orange"], badge, Pt(10))


def create_process_diagram(slide):
    """프로세스 다이어그램"""
    steps = [
        ("1", "수요 취합", "수주+재고"),
        ("2", "자원 확인", "설비/인력"),
        ("3", "최적화", "스케줄링"),
        ("4", "확정", "승인→지시"),
    ]
    
    start_x = Inches(0.5)
    y = Inches(3)
    box_width = Inches(2)
    box_height = Inches(1.5)
    gap = Inches(0.4)
    
    for i, (num, name, desc) in enumerate(steps):
        x = start_x + i * (box_width + gap)
        
        # 번호 원
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + box_width/2 - Inches(0.3), 
                                        y - Inches(0.5), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS["blue"]
        circle.line.fill.background()
        
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]
        p.alignment = PP_ALIGN.CENTER
        
        # 박스
        box = add_rounded_rect(slide, x, y, box_width, box_height, COLORS["light_gray"])
        
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS["navy"]
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS["gray"]
        p.alignment = PP_ALIGN.CENTER
        
        # 화살표 (마지막 제외)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, 
                                           x + box_width, y + box_height/2 - Inches(0.15),
                                           gap - Inches(0.1), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS["blue"]
            arrow.line.fill.background()


def create_comparison_diagram(slide):
    """Before/After 비교 다이어그램"""
    # Before 박스
    before_box = add_rounded_rect(slide, Inches(0.5), Inches(2.5), Inches(4), Inches(3), COLORS["light_gray"])
    tf = before_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "AS-IS"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS["gray"]
    p.alignment = PP_ALIGN.CENTER
    
    for item in ["✗ Excel 수작업", "✗ 제약조건 미반영", "✗ 납기율 저조"]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["dark"]
        p.alignment = PP_ALIGN.CENTER
    
    # 화살표
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.5), Inches(3.5), Inches(1), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLORS["orange"]
    arrow.line.fill.background()
    
    # After 박스
    after_box = add_rounded_rect(slide, Inches(5.5), Inches(2.5), Inches(4), Inches(3), COLORS["blue"])
    tf = after_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "TO-BE"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER
    
    for item in ["✓ 시스템 자동화", "✓ OR-Tools 최적화", "✓ 납기율 95%+"]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["white"]
        p.alignment = PP_ALIGN.CENTER


def create_conclusion_slide(slide):
    """결론 슬라이드"""
    # 중앙 강조 박스
    main_box = add_rounded_rect(slide, Inches(1), Inches(2), Inches(8), Inches(2), COLORS["navy"])
    tf = main_box.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "이제 우리는 이렇게 일해야 합니다"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = '"변화는 선택이 아닌 생존입니다"'
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = COLORS["light_blue"]
    p.alignment = PP_ALIGN.CENTER
    
    # Before/After 미니 박스
    add_rounded_rect(slide, Inches(1), Inches(4.5), Inches(3.8), Inches(2), COLORS["gray"],
                    "Before\n\n경험 의존 / Excel 수기 / 월말 정산", Pt(14))
    
    add_rounded_rect(slide, Inches(5.2), Inches(4.5), Inches(3.8), Inches(2), COLORS["green"],
                    "After\n\n데이터 기반 / 자동화 / 실시간 파악", Pt(14))


def main():
    print("📂 PPT 파일 로드 중...")
    prs = Presentation(INPUT_FILE)
    print(f"   총 {len(prs.slides)}개 슬라이드")
    
    # 슬라이드별 다이어그램 추가
    slide_actions = {
        0: ("표지", create_slide_1_title),
        1: ("목차", create_slide_2_toc),
        6: ("MES 구성도", create_mes_diagram),
        7: ("MES 구성도", create_mes_diagram),
        8: ("로드맵", create_roadmap_diagram),
        9: ("Before/After", create_comparison_diagram),
        10: ("프로세스", create_process_diagram),
        13: ("결론", create_conclusion_slide),
    }
    
    for idx, (name, func) in slide_actions.items():
        if idx < len(prs.slides):
            print(f"\n📊 슬라이드 {idx+1}: {name} 다이어그램 추가")
            func(prs.slides[idx])
            print("   ✅ 완료")
    
    # 저장
    print(f"\n💾 저장 중: {OUTPUT_FILE}")
    prs.save(OUTPUT_FILE)
    print("✅ 최종본 PPT 생성 완료!")
    print(f"   📁 {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
