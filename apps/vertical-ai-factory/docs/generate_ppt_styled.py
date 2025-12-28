"""
우수 제안서 스타일 PPT 생성 스크립트
- 2024 트렌드 반영: 미니멀리즘, 네이비 계열, Pretendard 폰트 스타일
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# 파일 경로
INPUT_FILE = '/home/roarm_m3/ai-factory-lab/apps/vertical-ai-factory/docs/생산계획시스템_설명자료_260108.pptx'
OUTPUT_FILE = '/home/roarm_m3/ai-factory-lab/apps/vertical-ai-factory/docs/생산계획시스템_설명자료_프로.pptx'

# =============================================================================
# 스타일 가이드 (2024 우수 제안서 트렌드)
# =============================================================================

STYLE = {
    # 색상 팔레트
    "primary": RGBColor(0x1E, 0x3A, 0x5F),      # 네이비 (신뢰감)
    "secondary": RGBColor(0x3B, 0x82, 0xF6),    # 파란색 (강조)
    "accent": RGBColor(0xF5, 0x9E, 0x0B),       # 주황 (포인트)
    "text_dark": RGBColor(0x1F, 0x29, 0x37),    # 진한 텍스트
    "text_light": RGBColor(0xFF, 0xFF, 0xFF),   # 흰색 텍스트
    "text_gray": RGBColor(0x6B, 0x72, 0x80),    # 회색 텍스트
    "bg_light": RGBColor(0xF8, 0xFA, 0xFC),     # 연한 배경
    
    # 폰트 (한글 환경)
    "font_title": "맑은 고딕",
    "font_body": "맑은 고딕",
    
    # 크기
    "title_size": Pt(32),
    "subtitle_size": Pt(18),
    "body_size": Pt(14),
    "small_size": Pt(11),
}


# =============================================================================
# 슬라이드 컨텐츠 (스타일 적용 버전)
# =============================================================================

SLIDES = {
    1: {
        "type": "title",
        "title": "차세대 생산관리 시스템",
        "subtitle": "디지털 전환을 통한 제조 혁신",
        "date": "2025.01"
    },
    2: {
        "type": "toc",
        "title": "목차",
        "items": [
            ("01", "Intro", "왜 이 이야기를 하는가"),
            ("02", "서론", "시대는 이미 변했다"),
            ("03", "본론", "무엇을 어떻게 바꿀 것인가"),
            ("04", "결론", "일하는 방식의 변화"),
        ]
    },
    3: {
        "type": "section",
        "number": "01",
        "title": "인사말 & 회사 소개"
    },
    4: {
        "type": "content",
        "section": "01 Intro",
        "title": "회사 소개",
        "content": [
            {"icon": "🏢", "title": "바이너리소프트", "desc": "20년 이상 ERP/MES 전문"},
            {"icon": "🏆", "title": "스마트팩토리 A등급", "desc": "검증된 공급기업"},
            {"icon": "🔧", "title": "핵심 역량", "desc": "MES 구축, AI 최적화, 실시간 모니터링"},
        ]
    },
    5: {
        "type": "content",
        "section": "02 서론",
        "title": "시대 흐름의 변화",
        "columns": [
            {"label": "과거", "color": "gray", "items": ["대량생산, 소품종", "Excel 수작업", "월 단위 집계"]},
            {"label": "현재", "color": "blue", "items": ["다품종 소량생산", "AI 수요예측", "실시간 수집"]},
            {"label": "미래", "color": "accent", "items": ["자율 생산", "디지털 트윈", "탄소중립"]},
        ]
    },
    6: {
        "type": "content",
        "section": "02 서론",
        "title": "어떻게 준비해야 하는가",
        "numbered": [
            {"num": "1", "title": "데이터 기반 의사결정", "desc": "실시간 생산 현황, KPI 대시보드"},
            {"num": "2", "title": "프로세스 표준화", "desc": "생산계획 절차 정립, 실적 자동화"},
            {"num": "3", "title": "시스템 통합", "desc": "ERP ↔ MES 연계, IoT 수집"},
        ]
    },
    7: {
        "type": "section",
        "number": "02",
        "title": "차세대 생산관리 시스템"
    },
    8: {
        "type": "diagram",
        "section": "03 본론",
        "title": "MES 전체 구성도",
        "layers": [
            ("ERP", "계획 / 관리", "primary"),
            ("MES", "실행 / 모니터링 / 추적", "secondary"),
            ("설비층", "PLC / 센서 / SCADA", "gray"),
        ]
    },
    9: {
        "type": "roadmap",
        "section": "03 본론",
        "title": "단계별 구축 전략",
        "phases": [
            {"phase": "1단계", "name": "생산계획", "status": "현재 목표", "items": ["주간/일간 계획", "OR-Tools 최적화", "작업지시 자동화"]},
            {"phase": "2단계", "name": "실적관리", "status": "확장", "items": ["바코드 실적 수집", "실시간 진척 현황"]},
            {"phase": "3단계", "name": "원가관리", "status": "고도화", "items": ["제조원가 집계", "원가 시뮬레이션"]},
        ]
    },
    10: {
        "type": "comparison",
        "section": "03 본론",
        "title": "생산계획 혁신 방향",
        "before": ["Excel 수작업 계획", "제약조건 미반영", "납기 준수율 저조"],
        "after": ["시스템 자동 계획", "OR-Tools 최적화", "납기 95% 이상"],
        "effect": "계획 수립 시간 80% 단축"
    },
    11: {
        "type": "process",
        "section": "03 본론",
        "title": "생산계획 프로세스",
        "steps": [
            {"num": "1", "name": "수요 취합", "desc": "수주 + 재고계획"},
            {"num": "2", "name": "자원 확인", "desc": "설비, 인력, 자재"},
            {"num": "3", "name": "최적화", "desc": "OR-Tools 스케줄링"},
            {"num": "4", "name": "계획 확정", "desc": "검토 → 승인 → 지시"},
        ]
    },
    12: {
        "type": "section",
        "number": "03",
        "title": "결론"
    },
    13: {
        "type": "expansion",
        "section": "03 확장",
        "title": "지속적인 고도화",
        "items": [
            {"title": "실적 관리", "bullets": ["바코드/터치 입력", "실시간 모니터링", "불량/재작업 추적"]},
            {"title": "원가 관리", "bullets": ["제조원가 집계", "직재/노무/경비 배부", "원가 시뮬레이션"]},
            {"title": "PO 도입", "bullets": ["생산오더 기반 관리", "오더별 추적/이력"]},
        ]
    },
    14: {
        "type": "conclusion",
        "title": "이제 우리는\n이렇게 일해야 합니다",
        "quote": '"변화는 선택이 아닌 생존입니다"',
        "items": [
            ("Before", ["경험과 감 의존", "Excel 수기", "월말 정산"]),
            ("After", ["데이터 기반", "시스템 자동화", "실시간 파악"]),
        ]
    }
}


def add_styled_text_box(slide, left, top, width, height, text, font_size, font_color, bold=False, align=PP_ALIGN.LEFT):
    """스타일이 적용된 텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = STYLE["font_body"]
    p.alignment = align
    return txBox


def add_rectangle(slide, left, top, width, height, fill_color, line_color=None):
    """색상이 채워진 사각형 추가"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def create_title_slide(prs, slide_idx, data):
    """표지 슬라이드 스타일링"""
    slide = prs.slides[slide_idx]
    
    # 기존 shapes 제거하지 않고 텍스트만 업데이트
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            if '차세대' in shape.text or '생산관리' in shape.text or '제안' in shape.text:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = data["title"]
                p.font.size = Pt(44)
                p.font.bold = True
                p.font.color.rgb = STYLE["primary"]
                p.font.name = STYLE["font_title"]
                
                p2 = tf.add_paragraph()
                p2.text = data["subtitle"]
                p2.font.size = Pt(20)
                p2.font.color.rgb = STYLE["text_gray"]
                p2.font.name = STYLE["font_body"]
            
            elif '2025' in shape.text or '2024' in shape.text:
                shape.text_frame.paragraphs[0].text = data["date"]
                shape.text_frame.paragraphs[0].font.color.rgb = STYLE["secondary"]


def create_section_slide(prs, slide_idx, data):
    """섹션 구분 슬라이드"""
    slide = prs.slides[slide_idx]
    
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text.strip()
            if any(x in text for x in ['제안', '인트로', '내용', '결론']):
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = data["number"]
                p.font.size = Pt(72)
                p.font.bold = True
                p.font.color.rgb = STYLE["secondary"]
            elif len(text) > 5:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = data["title"]
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = STYLE["primary"]


def create_content_slide(prs, slide_idx, data):
    """내용 슬라이드 스타일링"""
    slide = prs.slides[slide_idx]
    
    text_boxes = [s for s in slide.shapes if hasattr(s, 'text_frame')]
    
    if len(text_boxes) >= 2:
        # 섹션 번호
        tf = text_boxes[0].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = data.get("section", "")
        p.font.size = Pt(14)
        p.font.color.rgb = STYLE["secondary"]
        
        # 제목 + 내용
        tf = text_boxes[1].text_frame
        tf.clear()
        
        # 제목
        p = tf.paragraphs[0]
        p.text = data["title"]
        p.font.size = STYLE["title_size"]
        p.font.bold = True
        p.font.color.rgb = STYLE["primary"]
        
        # 내용 (유형에 따라)
        if "content" in data:
            for item in data["content"]:
                p = tf.add_paragraph()
                p.text = f"{item['icon']}  {item['title']}"
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = STYLE["text_dark"]
                
                p = tf.add_paragraph()
                p.text = f"      {item['desc']}"
                p.font.size = Pt(14)
                p.font.color.rgb = STYLE["text_gray"]
                p.space_after = Pt(12)
                
        elif "numbered" in data:
            for item in data["numbered"]:
                p = tf.add_paragraph()
                p.text = ""
                p.space_before = Pt(16)
                
                p = tf.add_paragraph()
                p.text = f"  {item['num']}   {item['title']}"
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = STYLE["secondary"]
                
                p = tf.add_paragraph()
                p.text = f"        {item['desc']}"
                p.font.size = Pt(14)
                p.font.color.rgb = STYLE["text_gray"]
                
        elif "columns" in data:
            for col in data["columns"]:
                p = tf.add_paragraph()
                p.text = f"\n【{col['label']}】"
                p.font.size = Pt(18)
                p.font.bold = True
                if col["color"] == "accent":
                    p.font.color.rgb = STYLE["accent"]
                elif col["color"] == "blue":
                    p.font.color.rgb = STYLE["secondary"]
                else:
                    p.font.color.rgb = STYLE["text_gray"]
                
                for item in col["items"]:
                    p = tf.add_paragraph()
                    p.text = f"    • {item}"
                    p.font.size = Pt(14)
                    p.font.color.rgb = STYLE["text_dark"]


def create_conclusion_slide(prs, slide_idx, data):
    """결론 슬라이드"""
    slide = prs.slides[slide_idx]
    
    text_boxes = [s for s in slide.shapes if hasattr(s, 'text_frame')]
    
    if len(text_boxes) >= 2:
        # 제목
        tf = text_boxes[1].text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = data["title"]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = STYLE["primary"]
        p.alignment = PP_ALIGN.CENTER
        
        # 인용문
        p = tf.add_paragraph()
        p.text = ""
        p = tf.add_paragraph()
        p.text = data["quote"]
        p.font.size = Pt(24)
        p.font.italic = True
        p.font.color.rgb = STYLE["secondary"]
        p.alignment = PP_ALIGN.CENTER
        
        # Before / After
        p = tf.add_paragraph()
        p.text = ""
        
        for label, items in data["items"]:
            p = tf.add_paragraph()
            p.text = f"\n【{label}】"
            p.font.size = Pt(18)
            p.font.bold = True
            if label == "After":
                p.font.color.rgb = STYLE["accent"]
            else:
                p.font.color.rgb = STYLE["text_gray"]
            
            for item in items:
                p = tf.add_paragraph()
                p.text = f"  • {item}"
                p.font.size = Pt(14)
                p.font.color.rgb = STYLE["text_dark"]


def main():
    print("📂 PPT 파일 로드 중...")
    prs = Presentation(INPUT_FILE)
    print(f"   총 {len(prs.slides)}개 슬라이드")
    
    # 각 슬라이드 스타일 적용
    for slide_idx, (slide_num, data) in enumerate(SLIDES.items()):
        if slide_num > len(prs.slides):
            continue
            
        print(f"\n🎨 슬라이드 {slide_num} 스타일링: {data.get('title', data.get('type', ''))}")
        
        slide_type = data.get("type", "content")
        
        if slide_type == "title":
            create_title_slide(prs, slide_idx, data)
        elif slide_type == "section":
            create_section_slide(prs, slide_idx, data)
        elif slide_type == "conclusion":
            create_conclusion_slide(prs, slide_idx, data)
        else:
            create_content_slide(prs, slide_idx, data)
        
        print("   ✅ 완료")
    
    # 저장
    print(f"\n💾 저장 중: {OUTPUT_FILE}")
    prs.save(OUTPUT_FILE)
    print("✅ 프로페셔널 PPT 생성 완료!")
    print(f"   📁 {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
