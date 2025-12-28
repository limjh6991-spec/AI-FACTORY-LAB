"""
생산계획시스템 PPT 자동 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy
import os

# 파일 경로
INPUT_FILE = '/home/roarm_m3/ai-factory-lab/apps/vertical-ai-factory/docs/생산계획시스템_설명자료_260108.pptx'
OUTPUT_FILE = '/home/roarm_m3/ai-factory-lab/apps/vertical-ai-factory/docs/생산계획시스템_설명자료_완성본.pptx'

def update_text_in_slide(slide, old_text_part, new_text):
    """슬라이드에서 특정 텍스트를 찾아 교체"""
    for shape in slide.shapes:
        if hasattr(shape, 'text') and old_text_part in shape.text:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if old_text_part in run.text:
                            run.text = run.text.replace(old_text_part, new_text)
                            return True
    return False

def set_slide_content(slide, title_text, content_text=None):
    """슬라이드 제목과 내용 설정"""
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text.strip()
            # 첫 번째 텍스트 박스를 제목으로 간주
            if text and len(text) < 50:
                shape.text_frame.paragraphs[0].runs[0].text = title_text
                break

def add_content_to_slide(slide, content_lines):
    """슬라이드에 내용 추가 (기존 내용 아래에)"""
    # 내용을 담을 텍스트박스 찾기 또는 생성
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            tf = shape.text_frame
            if len(tf.paragraphs) > 0:
                # 기존 설명 텍스트를 새 내용으로 교체
                text = shape.text
                if '페이지' in text or 'image' in text or '상세' in text:
                    tf.clear()
                    for i, line in enumerate(content_lines):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = line
                        p.font.size = Pt(14)
                    return True
    return False


# =============================================================================
# 슬라이드별 컨텐츠 정의
# =============================================================================

SLIDE_CONTENTS = {
    # 슬라이드 3: 인사말 & 회사 소개
    3: {
        "section": "01. Intro",
        "title": "인사말 & 회사 소개",
        "bullets": [
            "• 안녕하세요, 바이너리소프트입니다",
            "• 20년 이상의 ERP/MES 구축 경험",
            "• 스마트팩토리 A등급 공급기업",
            "",
            "◆ 핵심 역량",
            "  - 제조실행시스템(MES) 구축",
            "  - 생산계획 최적화 (OR-Tools, AI)",
            "  - 실시간 모니터링 솔루션",
            "",
            "◆ 주요 고객사",
            "  - 자동차부품, 전자, 기계 업종 다수"
        ]
    },
    
    # 슬라이드 4: 시대 흐름의 변화
    4: {
        "section": "01. Intro",
        "title": "시대 흐름의 변화",
        "bullets": [
            "◆ 제조업 환경의 변화",
            "",
            "【과거】",
            "  • 대량생산, 소품종",
            "  • 수작업 생산계획 (Excel)",
            "  • 월 단위 실적 집계",
            "",
            "【현재】",
            "  • 다품종 소량생산",
            "  • AI 기반 수요예측",
            "  • 실시간 데이터 수집",
            "",
            "【미래】",
            "  • 자율 생산 (Autonomous)",
            "  • 디지털 트윈 기반 시뮬레이션",
            "  • 탄소중립 대응 생산관리"
        ]
    },
    
    # 슬라이드 5: 우리는 어떻게 준비해야 하는가?
    5: {
        "section": "02. 서론",
        "title": "우리는 어떻게 준비해야 하는가?",
        "bullets": [
            "◆ 디지털 전환의 핵심 3요소",
            "",
            "1️⃣ 데이터 기반 의사결정",
            "   • 실시간 생산 현황 파악",
            "   • KPI 대시보드 구축",
            "",
            "2️⃣ 프로세스 표준화",
            "   • 생산계획 수립 절차 정립",
            "   • 실적 집계 자동화",
            "",
            "3️⃣ 시스템 통합",
            "   • ERP ↔ MES 연계",
            "   • 설비 데이터 수집 (PLC/IoT)"
        ]
    },
    
    # 슬라이드 7: MES 개요와 전체 구성도
    7: {
        "section": "02. 본론",
        "title": "MES 개요와 전체 구성도",
        "bullets": [
            "◆ MES (Manufacturing Execution System)",
            "",
            "┌─────────────────────────────────┐",
            "│          ERP (계획/관리)         │",
            "├─────────────────────────────────┤",
            "│   MES (실행/모니터링/추적)       │",
            "│  ┌─────┬─────┬─────┬─────┐    │",
            "│  │생산  │품질 │설비  │재고  │    │",
            "│  │관리  │관리 │관리  │관리  │    │",
            "│  └─────┴─────┴─────┴─────┘    │",
            "├─────────────────────────────────┤",
            "│     설비층 (PLC/센서/SCADA)      │",
            "└─────────────────────────────────┘"
        ]
    },
    
    # 슬라이드 8: 우리의 현실에 맞는 MES 구성도
    8: {
        "section": "02. 본론",
        "title": "우리의 현실에 맞는 MES 구성도",
        "bullets": [
            "◆ 단계별 구축 전략 (Phase Approach)",
            "",
            "【1단계: 생산계획】 ← 현재 목표",
            "  • 주간/일간 생산계획 수립",
            "  • OR-Tools 기반 최적화",
            "  • 작업지시서 자동 생성",
            "",
            "【2단계: 실적관리】",
            "  • 바코드/RFID 실적 수집",
            "  • 실시간 진척 현황",
            "",
            "【3단계: 원가관리】",
            "  • 제조원가 자동 집계",
            "  • 원가 시뮬레이션"
        ]
    },
    
    # 슬라이드 9: 생산계획 혁신 방향
    9: {
        "section": "02. 본론",
        "title": "생산계획 혁신 방향",
        "bullets": [
            "◆ AS-IS → TO-BE",
            "",
            "【현재 문제점】",
            "  ✗ Excel 수작업 계획",
            "  ✗ 설비/인력 제약 미반영",
            "  ✗ 납기 준수율 저조",
            "",
            "【개선 방향】",
            "  ✓ 시스템 기반 자동 계획",
            "  ✓ 제약조건 최적화 (OR-Tools)",
            "  ✓ 실시간 계획 변경 대응",
            "",
            "【기대 효과】",
            "  • 계획 수립 시간 80% 단축",
            "  • 납기 준수율 95% 이상"
        ]
    },
    
    # 슬라이드 10: 생산계획 프로세스 상세
    10: {
        "section": "02. 본론",
        "title": "생산계획 프로세스 상세",
        "bullets": [
            "◆ 생산계획 수립 프로세스",
            "",
            "┌──────────────────────────────────┐",
            "│  1. 수요 취합                      │",
            "│     └─ 수주 + 재고계획 데이터      │",
            "├──────────────────────────────────┤",
            "│  2. 자원 확인                      │",
            "│     └─ 설비능력, 인력, 자재        │",
            "├──────────────────────────────────┤",
            "│  3. 최적화 엔진                    │",
            "│     └─ OR-Tools Job Shop 스케줄링 │",
            "├──────────────────────────────────┤",
            "│  4. 계획 확정                      │",
            "│     └─ 검토 → 승인 → 작업지시     │",
            "└──────────────────────────────────┘"
        ]
    },
    
    # 슬라이드 12: 지속적인 고도화
    12: {
        "section": "03. 확장",
        "title": "지속적인 고도화 방향",
        "bullets": [
            "◆ 1단계 이후 확장 로드맵",
            "",
            "【실적 관리】",
            "  • 바코드/터치 기반 실적 입력",
            "  • 실시간 생산 현황 모니터링",
            "  • 불량/재작업 추적",
            "",
            "【원가 관리】",
            "  • 작업별 제조원가 집계",
            "  • 직접재료비/노무비/경비 배부",
            "  • 원가 절감 시뮬레이션",
            "",
            "【PO 개념 도입】",
            "  • 생산오더 기반 관리",
            "  • 오더별 추적/이력 관리"
        ]
    },
    
    # 슬라이드 13: 결론
    13: {
        "section": "03. 결론",
        "title": "이제 우리는 이렇게 일해야 합니다",
        "bullets": [
            "◆ 일하는 방식의 변화",
            "",
            "【Before】",
            "  • 경험과 감에 의존",
            "  • Excel / 수기 관리",
            "  • 월말 정산",
            "",
            "【After】",
            "  • 데이터 기반 의사결정",
            "  • 시스템 자동화",
            "  • 실시간 현황 파악",
            "",
            "━━━━━━━━━━━━━━━━━━━━━━",
            "\"변화는 선택이 아닌 생존입니다\"",
            "━━━━━━━━━━━━━━━━━━━━━━"
        ]
    }
}


def main():
    print("📂 PPT 파일 로드 중...")
    prs = Presentation(INPUT_FILE)
    
    print(f"   총 {len(prs.slides)}개 슬라이드")
    
    # 각 슬라이드 내용 업데이트
    for slide_num, content in SLIDE_CONTENTS.items():
        print(f"\n📝 슬라이드 {slide_num} 업데이트: {content['title']}")
        
        if slide_num <= len(prs.slides):
            slide = prs.slides[slide_num - 1]
            
            # 기존 텍스트 박스들 확인 및 업데이트
            text_boxes = []
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    text_boxes.append(shape)
            
            # 일반적으로 첫 번째는 섹션, 두 번째는 제목/내용
            if len(text_boxes) >= 2:
                # 두 번째 텍스트 박스에 내용 채우기
                tf = text_boxes[1].text_frame
                tf.clear()
                
                # 제목
                p = tf.paragraphs[0]
                p.text = content['title']
                p.font.size = Pt(28)
                p.font.bold = True
                
                # 내용
                for line in content['bullets']:
                    p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(14)
                    p.level = 0
                
                print(f"   ✅ 완료")
            else:
                print(f"   ⚠️ 텍스트 박스 부족 ({len(text_boxes)}개)")
    
    # 저장
    print(f"\n💾 저장 중: {OUTPUT_FILE}")
    prs.save(OUTPUT_FILE)
    print("✅ PPT 생성 완료!")
    print(f"   📁 {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
